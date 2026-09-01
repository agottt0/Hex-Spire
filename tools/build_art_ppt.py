# -*- coding: utf-8 -*-
"""
Build the art high-concept deck for Hex Spire from docs/01_美术高概念.md.

Re-runnable:  python tools/build_art_ppt.py
Output:       docs/Hex Spire 美术高概念 v0.1.pptx

Content is transcribed from docs/01_美术高概念.md v0.1. When that doc changes,
update the DATA blocks below and re-run — do not hand-edit the .pptx.
"""
import os
import shutil
import tempfile

from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ART = os.path.join(ROOT, "Art")
OUT = os.path.join(ROOT, "docs", "Hex Spire 美术高概念 v0.1.pptx")
TMP = tempfile.mkdtemp(prefix="hexspire_ppt_")

# ---------------------------------------------------------------- palette
# from doc §4.1 (scene) and §4.4 (functional / UI-only)
INK = RGBColor(0x0B, 0x0C, 0x17)
NIGHT = RGBColor(0x13, 0x1A, 0x23)
STEEL = RGBColor(0x2D, 0x39, 0x49)
INDIGO = RGBColor(0x3D, 0x2B, 0x45)
LILAC = RGBColor(0x6C, 0x44, 0x64)
ASH = RGBColor(0xEE, 0xE7, 0xDD)
VERM = RGBColor(0x6D, 0x04, 0x09)
LANTERN = RGBColor(0x8A, 0x2B, 0x1C)

FN_DANGER = RGBColor(0xFF, 0x3B, 0x30)
FN_ALLY = RGBColor(0x3D, 0xD6, 0xD0)
FN_BLOCK = RGBColor(0x7F, 0xA8, 0xD9)
FN_GOLD = RGBColor(0xFF, 0xC9, 0x4D)
FN_RUNE = RGBColor(0xB4, 0x7F, 0xE8)
FN_HEAL = RGBColor(0x6E, 0xD8, 0x9A)

MUTED = RGBColor(0x8A, 0x93, 0xA6)
DIM = RGBColor(0x5B, 0x64, 0x76)
HAIR = RGBColor(0x26, 0x2E, 0x3B)
PANEL = RGBColor(0x11, 0x17, 0x1F)
PANEL2 = RGBColor(0x18, 0x20, 0x2B)

SANS = "微软雅黑"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
ML = Inches(0.72)
CW = Inches(11.893)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
_page = {"n": 0}
_imgn = {"n": 0}


# ------------------------------------------------------------- primitives
def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, int(l), int(t), int(w), int(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def txt(s, l, t, w, h, runs, size=14, color=ASH, bold=False, font=SANS,
        align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP, space_after=0):
    """runs = str | [para, ...]; para = str | [(text, {overrides}), ...]"""
    tb = s.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paras = runs if isinstance(runs, list) else [runs]
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(space_after)
        for text, spec in (p if isinstance(p, list) else [(p, {})]):
            r = para.add_run()
            r.text = text
            f = r.font
            f.name = spec.get("font", font)
            f.size = Pt(spec.get("size", size))
            f.bold = spec.get("bold", bold)
            f.italic = spec.get("italic", False)
            f.color.rgb = spec.get("color", color)
    return tb


def pic(s, name, l, t, w, h, darken=None, tint=None, fade_left=None, fade_bottom=None,
        crop_top=None, focus=0.5):
    path = os.path.join(ART, name)
    im = Image.open(path).convert("RGB")
    if crop_top:
        iw, ih = im.size
        im = im.crop((0, int(ih * crop_top[0]), iw, int(ih * crop_top[1])))
    tw, th = float(w), float(h)
    target = tw / th
    iw, ih = im.size
    if iw / ih > target:
        nw = int(ih * target)
        x0 = int((iw - nw) * focus)
        im = im.crop((x0, 0, x0 + nw, ih))
    else:
        nh = int(iw / target)
        y0 = int((ih - nh) * focus)
        im = im.crop((0, y0, iw, y0 + nh))
    if darken:
        im = ImageEnhance.Brightness(im).enhance(1 - darken)
    if tint:
        im = Image.blend(im, Image.new("RGB", im.size, (0x0B, 0x0C, 0x17)), tint)
    if fade_left:
        base = Image.new("RGB", im.size, (0x0B, 0x0C, 0x17))
        mask = Image.new("L", im.size, 0)
        px = mask.load()
        wpx, hpx = im.size
        edge = max(2, int(wpx * fade_left))
        for x in range(edge):
            v = 255 if x < edge * 0.5 else int(255 * (edge - x) / (edge * 0.5))
            for y in range(hpx):
                px[x, y] = v
        im = Image.composite(base, im, mask)
    if fade_bottom:
        base = Image.new("RGB", im.size, (0x0B, 0x0C, 0x17))
        mask = Image.new("L", im.size, 0)
        px = mask.load()
        wpx, hpx = im.size
        edge = max(2, int(hpx * fade_bottom))
        for i in range(edge):
            y = hpx - 1 - i
            v = 255 if i < edge * 0.35 else int(255 * (edge - i) / (edge * 0.65))
            for x in range(wpx):
                px[x, y] = v
        im = Image.composite(base, im, mask)
    _imgn["n"] += 1
    fp = os.path.join(TMP, "img%03d.png" % _imgn["n"])
    im.save(fp, "PNG")
    return s.shapes.add_picture(fp, int(l), int(t), int(w), int(h))


def head(s, title, kicker=None, sub=None, subw=None, footer=True):
    _page["n"] += 1
    if kicker:
        txt(s, ML, Inches(0.40), Inches(9), Inches(0.26), kicker,
            size=11.5, color=LANTERN, bold=True, spacing=1.0)
    txt(s, ML, Inches(0.63), Inches(10.8), Inches(0.55), title,
        size=28, bold=True, color=ASH, spacing=1.0)
    rect(s, ML, Inches(1.26), Inches(0.60), Pt(3.2), fill=VERM)
    if sub:
        txt(s, ML, Inches(1.45), subw or Inches(11.6), Inches(0.4), sub,
            size=13, color=MUTED, spacing=1.35)
    txt(s, Inches(12.0), Inches(6.98), Inches(0.62), Inches(0.24), "%02d" % _page["n"],
        size=11, color=DIM, align=PP_ALIGN.RIGHT, spacing=1.0)
    if footer:
        txt(s, ML, Inches(6.98), Inches(6), Inches(0.24),
            "《（暂名）Hex Spire》 美术高概念 v0.1",
            size=9.5, color=RGBColor(0x39, 0x40, 0x4E), spacing=1.0)
    return s


def table(s, l, t, w, headers, rows, colw, rowh=Inches(0.36), hh=Inches(0.33),
          fsize=11.5, hsize=10.5, cell_colors=None, bold_col0=True):
    x = l
    rect(s, l, t, w, hh, fill=RGBColor(0x1D, 0x25, 0x32))
    for i, htext in enumerate(headers):
        txt(s, x + Inches(0.11), t + Inches(0.065), colw[i] - Inches(0.16), hh,
            htext, size=hsize, color=FN_GOLD, bold=True, spacing=1.0)
        x += colw[i]
    y = t + hh
    for ri, row in enumerate(rows):
        rh = rowh[ri] if isinstance(rowh, list) else rowh
        if ri % 2 == 0:
            rect(s, l, y, w, rh, fill=PANEL)
        rect(s, l, y, w, Pt(0.6), fill=HAIR)
        x = l
        for ci, cell in enumerate(row):
            c = ASH
            if cell_colors and (ri, ci) in cell_colors:
                c = cell_colors[(ri, ci)]
            txt(s, x + Inches(0.11), y + Inches(0.06), colw[ci] - Inches(0.16), rh,
                cell, size=fsize, color=c, bold=(bold_col0 and ci == 0), spacing=1.1)
            x += colw[ci]
        y += rh
    return y


def swatch(s, l, t, w, h, color, name, hexs, note=None):
    rect(s, l, t, w, h, fill=color, line=RGBColor(0x2C, 0x34, 0x42))
    txt(s, l, t + h + Inches(0.06), w + Inches(0.08), Inches(0.2), name, size=8.5,
        bold=True, color=ASH, spacing=1.0)
    txt(s, l, t + h + Inches(0.23), w + Inches(0.08), Inches(0.2), hexs, size=8,
        color=MUTED, font=MONO, spacing=1.0)
    if note:
        txt(s, l, t + h + Inches(0.39), w + Inches(0.08), Inches(0.2), note,
            size=8, color=DIM, spacing=1.0)


def hexa(s, cx, cy, size, fill=None, line=None, lw=1.25):
    """Pointy-top hexagon of width `size` (flat-to-flat) centred on (cx, cy).

    pptx HEXAGON is flat-top, so we rotate 90deg. Rotation happens about the
    shape's centre and does NOT change its bounding box, so the pre-rotation
    box must be the post-rotation box transposed: w = size * 2/sqrt(3) becomes
    the final height, and h = size becomes the final width.
    """
    w = size * 1.1547   # becomes the height after rotating
    h = size            # becomes the width after rotating
    sh = s.shapes.add_shape(MSO_SHAPE.HEXAGON,
                            int(cx - w / 2), int(cy - h / 2),
                            int(w), int(h))
    sh.rotation = 90
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = False
    return sh


def bullets(s, l, t, w, items, size=13, gap=0.33, color=ASH, marker="—",
            mcolor=VERM, spacing=1.28, numbered=False, mw=0.3):
    """marker: literal string for every item; numbered=True -> 1. 2. 3. ..."""
    y = t
    for i, it in enumerate(items):
        m = ("%d." % (i + 1)) if numbered else marker
        txt(s, l, y, Inches(mw), Inches(0.26), m, size=size, color=mcolor,
            bold=True, spacing=1.0)
        txt(s, l + Inches(mw), y, w - Inches(mw), Inches(0.3), it, size=size,
            color=color, spacing=spacing)
        y += Inches(gap)
    return y


def text_w(text, size):
    """Rough rendered width in EMU. CJK glyphs are full-width, ASCII ~0.55em."""
    em = 0.0
    for ch in text:
        em += 1.0 if ord(ch) > 0x2E7F else 0.55
    return Inches(em * size / 72.0)


def tag(s, l, t, text, color=FN_GOLD, size=10, h=Inches(0.26)):
    w = text_w(text, size) + Inches(0.26)
    rect(s, l, t, w, h, fill=None, line=color, lw=0.75)
    txt(s, l, t + Inches(0.035), w, Inches(0.2), text, size=size, color=color,
        bold=True, align=PP_ALIGN.CENTER, spacing=1.0)
    return l + w


def note(s, l, t, w, text, color=FN_GOLD, size=11.5, h=Inches(0.5)):
    rect(s, l, t, Pt(2.4), h, fill=color)
    txt(s, l + Inches(0.16), t, w - Inches(0.16), h, text, size=size, color=MUTED,
        spacing=1.32)


# =============================================================== 01 cover
s = slide()
pic(s, "45125e37736044174b3f5b.png", 0, 0, W, H, darken=0.34, tint=0.30)
rect(s, 0, 0, Inches(7.1), H, fill=INK)
pic(s, "45125e37736044174b3f5b.png", Inches(6.4), 0, Inches(6.933), H,
    darken=0.16, tint=0.10, fade_left=0.34, focus=0.72)
rect(s, 0, 0, W, Pt(4), fill=VERM)

txt(s, ML, Inches(1.30), Inches(6.0), Inches(0.3),
    "ART HIGH CONCEPT  ·  v0.1  ·  2026.08", size=12, color=LANTERN, bold=True,
    font=MONO, spacing=1.0)
txt(s, ML, Inches(1.72), Inches(6.4), Inches(1.5),
    [[("（暂名）", {"size": 30, "color": MUTED}), ("Hex Spire", {"size": 52})]],
    size=52, bold=True, color=ASH, spacing=1.0)
txt(s, ML, Inches(2.62), Inches(6.2), Inches(0.5), "美术高概念",
    size=33, bold=True, color=ASH, spacing=1.0)
rect(s, ML, Inches(3.30), Inches(0.9), Pt(3), fill=VERM)

txt(s, ML, Inches(3.62), Inches(5.7), Inches(0.9),
    "在熟悉的城市中，我们走进了不熟悉的世界。", size=17, color=LILAC, spacing=1.35)
txt(s, ML, Inches(4.20), Inches(5.9), Inches(1.0),
    [[("画面里 ", {}), ("90%", {"color": FN_ALLY, "bold": True}),
      (" 是你今天真的见过的城市，", {}), ("10%", {"color": FN_DANGER, "bold": True}),
      (" 不对劲。", {})],
     [("恐怖来自那 10%，不来自那 90%。", {"color": MUTED})]],
    size=14.5, color=ASH, spacing=1.4)

x = ML
for t_ in ["单人肉鸽", "六边形战棋", "卡牌构筑", "东方怪谈 × 现代城市"]:
    x = tag(s, x, Inches(5.30), t_, color=STEEL if t_ != "东方怪谈 × 现代城市" else LANTERN,
            size=10.5) + Inches(0.1)
txt(s, ML, Inches(5.86), Inches(6.2), Inches(0.6),
    [[("对齐 ", {}), ("00_系统策划案.md v0.3", {"font": MONO, "color": MUTED}),
      ("   ·   待决项 A1–A8 见末页", {})]],
    size=11, color=DIM, spacing=1.2)

# =============================================================== 02 agenda
s = slide()
head(s, "这份文档要回答什么", kicker="CONTENTS · 目录",
     sub="美术不是“画得好看”，而是“每一条都能被验收”。全篇 16 节，本 PPT 取其中最需要拍板与最容易翻车的部分。")

items = [
    ("01", "美术定位与三大支柱", "V3 可读性 > V1 可信度 > V2 克制", FN_ALLY),
    ("02", "核心转译公式", "怪谈是怎么长在城市里的（16 条词表）", LANTERN),
    ("03", "色彩系统", "场景色板 × 功能色，红色冲突的解法", VERM),
    ("04", "三章场景高概念", "黄泉线 / 百鬼市 / 阴司", LILAC),
    ("05", "角色与怪物", "四职业配比 · 1/3/6 格体型的视觉承载", FN_RUNE),
    ("06", "战场与 Tile 规范", "128px 尖顶六边形 · R1–R9 可读性硬规则", FN_BLOCK),
    ("07", "腐蚀度视觉表达", "一张底图 + 三组贴花 + 四档后处理", FN_DANGER),
    ("08", "红线 · M0 清单 · 待决项", "明确不做什么 + 怎么验收 + A1–A8", FN_GOLD),
]
y = Inches(2.18)
for i, (num, title, sub, c) in enumerate(items):
    col = i % 2
    row = i // 2
    l = ML + col * Inches(6.05)
    t = y + row * Inches(1.16)
    rect(s, l, t, Inches(5.65), Inches(0.98), fill=PANEL)
    rect(s, l, t, Pt(3), Inches(0.98), fill=c)
    txt(s, l + Inches(0.24), t + Inches(0.15), Inches(0.6), Inches(0.4), num,
        size=20, bold=True, color=c, font=MONO, spacing=1.0)
    txt(s, l + Inches(0.92), t + Inches(0.13), Inches(4.5), Inches(0.3), title,
        size=15, bold=True, color=ASH, spacing=1.0)
    txt(s, l + Inches(0.92), t + Inches(0.52), Inches(4.55), Inches(0.3), sub,
        size=11.5, color=MUTED, spacing=1.1)

# ======================================================= 03 L0 moodboard
s = slide()
head(s, "L0 总方向 Moodboard：已锁定，不再讨论", kicker="01 · 定位",
     sub="它已经完成了它的工作 —— 定基调、对外沟通、外包首页。后续所有资产以它为准，不再改它。")
pic(s, "总方向 Moodboard.png", ML, Inches(2.05), Inches(8.05), Inches(4.53),
    darken=0.04)
rect(s, ML, Inches(2.05), Inches(8.05), Inches(4.53), fill=None, line=HAIR, lw=1)

lx = Inches(9.05)
txt(s, lx, Inches(2.05), Inches(3.6), Inches(0.3), "L0 已锁定的五件事",
    size=13.5, bold=True, color=FN_GOLD, spacing=1.0)
bullets(s, lx, Inches(2.45), Inches(3.6), [
    "主色三元：夜蓝 / 靛紫 + 朱红 + 香灰白",
    "剪影 + 氛围光（背光 / 逆光为主）",
    "高对比、低中间调",
    "东方符号作为点缀，而非主体",
    "情绪四词：诡谲 / 压迫 / 孤独 / 神秘",
], size=12, gap=0.44, spacing=1.25)

txt(s, lx, Inches(4.75), Inches(3.6), Inches(0.3), "但 L0 不足以开工",
    size=13.5, bold=True, color=FN_DANGER, spacing=1.0)
txt(s, lx, Inches(5.12), Inches(3.62), Inches(1.4),
    [[("L0 解决的是“这个游戏长什么样”。", {})],
     [("开工还缺 ", {}), ("7 张 L1 moodboard", {"color": ASH, "bold": True}),
      ("（MB-01~07）。其中 ", {}),
      ("MB-03 战斗可读性", {"color": FN_ALLY, "bold": True}),
      (" 最容易被跳过，也最不该跳过 —— 其他 moodboard 教美术“怎么变好看”，只有它教美术“怎么别挡住玩家”。", {})]],
    size=11.5, color=MUTED, spacing=1.36)

# ================================================= 04 pillars
s = slide()
head(s, "三大美术支柱：优先级是硬的", kicker="01 · 支柱",
     sub="当氛围与可读性冲突时（一定会冲突），永远牺牲氛围。这是 45 分钟一局、要玩几百小时的肉鸽，不是一次性的恐怖体验。")

cols = [
    ("V1", "日常的可信度", INDIGO, LILAC,
     "城市部分必须写实、有生活垢、有年代感。城市不可信 → 异常不恐怖。",
     "任何一张场景图，遮住异常元素后，应该像一张普通的城市夜景照片。"),
    ("V2", "异常的克制", VERM, LANTERN,
     "每个画面的“不对劲”控制在 1–3 处，且必须是逻辑错误而非血浆。",
     "每张场景图能明确列出异常清单，条目数 ≤3。数不出来、或数出 8 条 = 不合格。"),
    ("V3", "战斗的可读性", RGBColor(0x14, 0x3A, 0x3A), FN_ALLY,
     "战斗中氛围让位于信息。玩家必须 0.5 秒内读出：谁在哪、占几格、朝哪、下回合要干什么。",
     "把画面降到 25% 缩略图，仍能分清己方 / 敌方 / 危害格 / 体型轮廓。"),
]
for i, (code, name, bg, acc, mean, crit) in enumerate(cols):
    l = ML + i * Inches(4.02)
    rect(s, l, Inches(2.3), Inches(3.78), Inches(3.55), fill=PANEL)
    rect(s, l, Inches(2.3), Inches(3.78), Pt(3.5), fill=acc)
    txt(s, l + Inches(0.26), Inches(2.5), Inches(1.2), Inches(0.5), code,
        size=30, bold=True, color=acc, font=MONO, spacing=1.0)
    txt(s, l + Inches(1.3), Inches(2.62), Inches(2.3), Inches(0.4), name,
        size=16, bold=True, color=ASH, spacing=1.0)
    txt(s, l + Inches(0.26), Inches(3.25), Inches(3.3), Inches(1.0), mean,
        size=12, color=ASH, spacing=1.35)
    rect(s, l + Inches(0.26), Inches(4.42), Inches(3.28), Pt(0.7), fill=HAIR)
    txt(s, l + Inches(0.26), Inches(4.58), Inches(3.3), Inches(0.24),
        "做完了怎么验收", size=10.5, bold=True, color=FN_GOLD, spacing=1.0)
    txt(s, l + Inches(0.26), Inches(4.86), Inches(3.3), Inches(0.9), crit,
        size=11.5, color=MUTED, spacing=1.34)

rect(s, ML, Inches(6.12), CW, Inches(0.62), fill=RGBColor(0x1E, 0x0A, 0x0C))
rect(s, ML, Inches(6.12), Pt(3.5), Inches(0.62), fill=FN_DANGER)
txt(s, ML + Inches(0.28), Inches(6.24), Inches(11.4), Inches(0.4),
    [[("优先级： ", {"color": MUTED}), ("V3 可读性", {"color": FN_ALLY, "bold": True}),
      ("  >  ", {"color": DIM}), ("V1 可信度", {"color": LILAC, "bold": True}),
      ("  >  ", {"color": DIM}), ("V2 克制", {"color": LANTERN, "bold": True}),
      ("      冲突时的处理不是“折中”，而是“牺牲氛围”。", {"color": MUTED})]],
    size=13.5, spacing=1.0)

# ================================================= 05 formula
s = slide()
head(s, "核心转译公式：怪谈是怎么长在城市里的", kicker="02 · 公式",
     sub="所有场景与怪物设计都走同一个公式。这是最需要被外包和 AI 生成理解的一节。")

fy = Inches(2.12)
FW = Inches(8.62)          # formula band stops short of the right-hand panel
rect(s, ML, fy, FW, Inches(0.95), fill=PANEL2)
parts = [("【传统母题】", LILAC, Inches(1.72)), ("【现代城市载体】", FN_BLOCK, Inches(2.08)),
         ("【一处日常逻辑错误】", FN_DANGER, Inches(2.48))]
x = ML + Inches(0.24)
txt(s, x, fy + Inches(0.3), Inches(1.2), Inches(0.4), "怪谈元素 =",
    size=14, bold=True, color=ASH, spacing=1.0)
x += Inches(1.06)
for i, (p, c, wdt) in enumerate(parts):
    if i:
        txt(s, x, fy + Inches(0.3), Inches(0.28), Inches(0.4), "×", size=15,
            bold=True, color=DIM, align=PP_ALIGN.CENTER, spacing=1.0)
        x += Inches(0.3)
    rect(s, x, fy + Inches(0.24), wdt, Inches(0.48), fill=None, line=c, lw=1.1)
    txt(s, x, fy + Inches(0.36), wdt, Inches(0.3), p, size=11.5, bold=True,
        color=c, align=PP_ALIGN.CENTER, spacing=1.0)
    x += wdt

miss = [("缺【传统母题】", "→ 变成普通都市废墟 / 克苏鲁，失去“东方”"),
        ("缺【现代城市载体】", "→ 变成古风灵异，失去“现代”，也失去玩家的代入"),
        ("缺【日常逻辑错误】", "→ 变成美术堆料，玩家只觉得“红红的挺好看”")]
for i, (a, b) in enumerate(miss):
    t = Inches(3.3) + i * Inches(0.42)
    txt(s, ML, t, Inches(2.4), Inches(0.3), a, size=12, bold=True, color=FN_DANGER, spacing=1.0)
    txt(s, ML + Inches(2.45), t, Inches(6.2), Inches(0.3), b, size=12, color=MUTED, spacing=1.0)

txt(s, ML, Inches(4.72), Inches(8.6), Inches(0.3), "转译词表节选（完整 16 条见文档 §2.1，持续扩充）",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
table(s, ML, Inches(5.08), Inches(8.62),
      ["传统母题", "现代城市载体", "日常逻辑错误（异常点）", "用在"],
      [["走无常", "末班车司机 / 安检员", "制服正确、工牌正确，但没有影子", "黄泉线"],
       ["纸扎（纸马纸屋）", "商场家电区", "冰箱电视是纸糊的，纸马停在扭蛋机旁", "百鬼市"],
       ["冥币面额", "价签 / 促销吊旗", "单位不是元，是“年”（用寿命付账）", "百鬼市"],
       ["孟婆汤", "医院配药窗口", "药袋上没有药名，只有一个“忘”字", "阴司"]],
      [Inches(1.7), Inches(2.1), Inches(3.7), Inches(1.12)],
      rowh=Inches(0.33), fsize=11)

rx = Inches(9.55)
rect(s, rx, Inches(2.12), Inches(3.06), Inches(4.62), fill=PANEL)
rect(s, rx, Inches(2.12), Inches(3.06), Pt(3), fill=FN_GOLD)
txt(s, rx + Inches(0.24), Inches(2.34), Inches(2.6), Inches(0.3), "使用规则",
    size=13.5, bold=True, color=FN_GOLD, spacing=1.0)
bullets(s, rx + Inches(0.24), Inches(2.74), Inches(2.62), [
    "一个房间最多用 3 条。用满 6 条不是恐怖，是杂货铺",
    "同一条在同一章内最多复用 2 次，第三次必须换载体",
    "每条都要能“遮住即消失”—— 异常必须是可移除的图层",
], size=11.5, gap=0.72, spacing=1.3, numbered=True, mcolor=FN_GOLD, mw=0.26)
rect(s, rx + Inches(0.24), Inches(4.98), Inches(2.6), Pt(0.7), fill=HAIR)
txt(s, rx + Inches(0.24), Inches(5.16), Inches(2.62), Inches(1.4),
    [[("第 3 点是省钱的关键：", {"color": FN_ALLY, "bold": True})],
     [("底图（城市）一张，异常图层按腐蚀度分 3 档叠加。同一个房间在腐蚀度 1 和 7 下看起来是两个地方，但美术只画了一张底图 + 三组贴花。", {})]],
    size=11, color=MUTED, spacing=1.34)

# ================================================= 06 palette
s = slide()
head(s, "全局色板：从 L0 实测采样，非估值", kicker="03 · 色彩",
     sub="每章换色相，不换明度结构。明度结构（暗底 + 极少高光）全局统一 —— 这是画面能一眼认出是同一个游戏的原因。")

sw = [(INK, "墨夜蓝 INK", "#0B0C17", "35–45%"), (NIGHT, "夜青 NIGHT", "#131A23", "15–20%"),
      (STEEL, "钢灰蓝 STEEL", "#2D3949", "15%"), (INDIGO, "靛紫 INDIGO", "#3D2B45", "8%"),
      (LILAC, "幽紫 LILAC", "#6C4464", "5%"), (ASH, "香灰白 ASH", "#EEE7DD", "4%"),
      (VERM, "朱红 VERM", "#6D0409", "≤3%"), (LANTERN, "灯笼橙红", "#8A2B1C", "≤3%")]
SWP = Inches(1.078)                       # pitch: 8 swatches inside 8.62in
for i, (c, n, hx, pc) in enumerate(sw):
    swatch(s, ML + i * SWP, Inches(2.15), Inches(0.98), Inches(0.78), c, n, hx, pc)
rect(s, ML + 6 * SWP, Inches(2.15), Inches(0.98), Inches(0.78),
     fill=None, line=FN_DANGER, lw=1.75)

txt(s, ML, Inches(3.86), Inches(8.5), Inches(0.3), "分章色板：换色相，不换明度结构",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
chs = [("第一章 · 地铁 · 黄泉线", "靛蓝 → 紫", "冷白荧光管 + 单点朱红应急灯", "又冷又空，只有你一个人",
        ["#0B0C17", "#292F4C", "#311726", "#5F4E73", "#B19CB8"]),
       ("第二章 · 商场 · 百鬼市", "暗红褐 → 暖橙", "大量灯笼暖光 + 冷蓝阴影", "热闹得不正常",
        ["#181012", "#391A1A", "#563C42", "#836065", "#C3A8AA"]),
       ("终章 · 医院 · 阴司", "青灰 → 冷绿", "全冷光、无暖色、无阴影", "干净、明亮、没有一点人味",
        ["#050D13", "#0B1E28", "#263B48", "#415D6F", "#88A4B4"])]
for i, (nm, hue, light, one, hexes) in enumerate(chs):
    t = Inches(4.22) + i * Inches(0.85)
    rect(s, ML, t, Inches(8.6), Inches(0.74), fill=PANEL if i % 2 == 0 else RGBColor(0x0E, 0x13, 0x1B))
    txt(s, ML + Inches(0.16), t + Inches(0.09), Inches(2.3), Inches(0.3), nm,
        size=12, bold=True, color=ASH, spacing=1.0)
    txt(s, ML + Inches(0.16), t + Inches(0.4), Inches(2.3), Inches(0.26), hue,
        size=10.5, color=LILAC, spacing=1.0)
    for j, hx in enumerate(hexes):
        c = RGBColor(int(hx[1:3], 16), int(hx[3:5], 16), int(hx[5:7], 16))
        rect(s, ML + Inches(2.6) + j * Inches(0.44), t + Inches(0.17),
             Inches(0.4), Inches(0.4), fill=c, line=RGBColor(0x2C, 0x34, 0x42))
    txt(s, ML + Inches(4.9), t + Inches(0.1), Inches(1.95), Inches(0.56), light,
        size=9.5, color=MUTED, spacing=1.2)
    txt(s, ML + Inches(6.95), t + Inches(0.1), Inches(1.55), Inches(0.56), one,
        size=9.5, color=FN_GOLD, spacing=1.2)

rx = Inches(9.55)
rect(s, rx, Inches(2.15), Inches(3.06), Inches(2.32), fill=RGBColor(0x1E, 0x0A, 0x0C))
rect(s, rx, Inches(2.15), Inches(3.06), Pt(3), fill=FN_DANGER)
txt(s, rx + Inches(0.22), Inches(2.36), Inches(2.66), Inches(0.3), "朱红 3% 上限是硬约束",
    size=12.5, bold=True, color=FN_DANGER, spacing=1.0)
txt(s, rx + Inches(0.22), Inches(2.74), Inches(2.66), Inches(1.6),
    [[("L0 里红看起来很多，因为它是一张宣传图（红在小图里抓眼）。实际场景里红一旦超过 5%，两件事同时发生：", {})],
     [("① 恐怖感变成廉价的血腥感", {"color": ASH})],
     [("② 红作为“伤害 / 危险”的 UI 功能色彻底失效", {"color": ASH})]],
    size=11, color=MUTED, spacing=1.34)

rect(s, rx, Inches(4.66), Inches(3.06), Inches(2.08), fill=PANEL)
rect(s, rx, Inches(4.66), Inches(3.06), Pt(3), fill=FN_ALLY)
txt(s, rx + Inches(0.22), Inches(4.86), Inches(2.66), Inches(0.3), "终章的反差设计",
    size=12.5, bold=True, color=FN_ALLY, spacing=1.0)
txt(s, rx + Inches(0.22), Inches(5.22), Inches(2.66), Inches(1.4),
    "前两章靠“暗”制造恐怖，终章靠“太亮太干净”。医院是唯一照明充足、几乎没有暗部、完全禁用暖色的章节。玩家在前 6 层学会“暗处有东西”，第三章突然没有暗处，反而更慌 —— 且不加一分资产量。",
    size=11, color=MUTED, spacing=1.34)

# ================================================= 07 red conflict
s = slide()
head(s, "必须先解决的冲突：红色", kicker="03 · 色彩 · 可读性地基",
     sub="本项目的氛围色是朱红（灯笼、符纸、安全线），而战棋的危险色也是红（伤害、敌方、非法落点）。二者同屏必然干扰。")

rect(s, ML, Inches(2.2), Inches(5.75), Inches(2.5), fill=RGBColor(0x14, 0x06, 0x08))
rect(s, ML, Inches(2.2), Inches(5.75), Pt(3), fill=VERM)
txt(s, ML + Inches(0.28), Inches(2.42), Inches(3.0), Inches(0.3), "场景红 · 装饰",
    size=13, bold=True, color=LANTERN, spacing=1.0)
rect(s, ML + Inches(0.28), Inches(2.82), Inches(1.5), Inches(1.0), fill=VERM,
     line=RGBColor(0x2C, 0x34, 0x42))
txt(s, ML + Inches(1.95), Inches(2.88), Inches(3.5), Inches(1.0),
    [[("#6D0409", {"font": MONO, "size": 15, "color": ASH, "bold": True})],
     [("色相 357°  ·  饱和 96%", {"color": MUTED})],
     [("明度 43%", {"color": FN_GOLD, "bold": True})]],
    size=11.5, color=MUTED, spacing=1.36)
txt(s, ML + Inches(0.28), Inches(4.0), Inches(5.2), Inches(0.55),
    "压暗、不动的红 = 装饰。灯笼、符纸、安全线、朱批。场景红明度上限 45%。",
    size=11.5, color=MUTED, spacing=1.34)

rect(s, ML, Inches(4.88), Inches(5.75), Inches(1.86), fill=PANEL)
rect(s, ML, Inches(4.88), Inches(5.75), Pt(3), fill=FN_DANGER)
txt(s, ML + Inches(0.28), Inches(5.08), Inches(3.0), Inches(0.3), "UI 红 · 威胁",
    size=13, bold=True, color=FN_DANGER, spacing=1.0)
rect(s, ML + Inches(0.28), Inches(5.46), Inches(1.5), Inches(0.7), fill=FN_DANGER,
     line=RGBColor(0x2C, 0x34, 0x42))
txt(s, ML + Inches(1.95), Inches(5.44), Inches(3.5), Inches(0.8),
    [[("#FF3B30", {"font": MONO, "size": 15, "color": ASH, "bold": True})],
     [("色相 3°  ·  饱和 81%  ·  ", {"color": MUTED}),
      ("明度 100%", {"color": FN_GOLD, "bold": True})]],
    size=11.5, color=MUTED, spacing=1.36)
txt(s, ML + Inches(0.28), Inches(6.28), Inches(5.2), Inches(0.35),
    "明亮、会动的红 = 威胁。伤害、敌方、非法落点、危害格。", size=11.5,
    color=MUTED, spacing=1.3)

rx = Inches(6.9)
txt(s, rx, Inches(2.2), Inches(5.7), Inches(0.66),
    [[("两者色相几乎相同，饱和度也都很高 —— 唯一的分野是“暗 vs 亮”。", {"color": ASH, "bold": True})],
     [("所以给外包的指令必须是「场景红要压暗（V≤45%），UI 红要拉亮（V=100%）」，不是「场景红要降饱和」。降饱和只会得到发灰的脏红，既不好看，也依然会与 UI 红混淆。", {})]],
    size=11.5, color=MUTED, spacing=1.36)

txt(s, rx, Inches(3.34), Inches(5.7), Inches(0.3), "功能色（UI 专用，任何场景不得占用）",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
fns = [(FN_DANGER, "FN_DANGER", "#FF3B30", "伤害 / 危险 / 敌方"),
       (FN_ALLY, "FN_ALLY", "#3DD6D0", "可移动 / 己方 / 合法"),
       (FN_BLOCK, "FN_BLOCK", "#7FA8D9", "格挡 / 防御（仅数值与图标）"),
       (FN_GOLD, "FN_GOLD", "#FFC94D", "暴击 / 稀有 / 奖励"),
       (FN_RUNE, "FN_RUNE", "#B47FE8", "符文 / 规则改写"),
       (FN_HEAL, "FN_HEAL", "#6ED89A", "治疗（仅数值与图标）")]
for i, (c, n, hx, use) in enumerate(fns):
    t = Inches(3.72) + i * Inches(0.42)
    rect(s, rx, t, Inches(0.34), Inches(0.32), fill=c)
    txt(s, rx + Inches(0.46), t + Inches(0.05), Inches(1.5), Inches(0.26), n,
        size=11, bold=True, color=ASH, font=MONO, spacing=1.0)
    txt(s, rx + Inches(1.95), t + Inches(0.05), Inches(0.95), Inches(0.26), hx,
        size=10.5, color=MUTED, font=MONO, spacing=1.0)
    txt(s, rx + Inches(2.95), t + Inches(0.05), Inches(2.7), Inches(0.26), use,
        size=10.5, color=MUTED, spacing=1.0)

rect(s, rx, Inches(6.3), Inches(5.7), Inches(0.44), fill=RGBColor(0x0A, 0x1F, 0x1F))
rect(s, rx, Inches(6.3), Pt(3), Inches(0.44), fill=FN_ALLY)
txt(s, rx + Inches(0.2), Inches(6.4), Inches(5.4), Inches(0.3),
    [[("地面高亮层只允许两色：", {"color": MUTED}), ("青蓝 = 合法/己方", {"color": FN_ALLY, "bold": True}),
      (" · ", {"color": DIM}), ("亮红 = 危险/非法", {"color": FN_DANGER, "bold": True}),
      ("。其余功能色只上 UI。", {"color": MUTED})]],
    size=11.5, spacing=1.0)

# ================================================= 08 chapter 1
s = slide()
_page["n"] += 1
pic(s, "地铁-黄泉线.png", 0, 0, W, H, darken=0.10, tint=0.06)
rect(s, 0, Inches(4.55), W, Inches(2.95), fill=INK)
pic(s, "地铁-黄泉线.png", 0, 0, W, Inches(5.1), darken=0.06, fade_bottom=0.30)
rect(s, 0, 0, W, Inches(1.15), fill=INK)
pic(s, "地铁-黄泉线.png", 0, 0, W, Inches(1.15), darken=0.62, tint=0.45, crop_top=(0.0, 0.16))
rect(s, 0, 0, Pt(4), H, fill=LILAC)
txt(s, ML, Inches(0.34), Inches(9), Inches(0.26), "04 · 场景 · 第一章",
    size=11.5, color=LILAC, bold=True, spacing=1.0)
txt(s, ML, Inches(0.58), Inches(9), Inches(0.45), "地铁 · 黄泉线",
    size=27, bold=True, color=ASH, spacing=1.0)
txt(s, Inches(6.3), Inches(0.66), Inches(6.3), Inches(0.4),
    "「通往未知的末班车，永远到不了终点。」", size=14.5, color=LILAC,
    align=PP_ALIGN.RIGHT, spacing=1.0)
txt(s, Inches(12.0), Inches(6.98), Inches(0.62), Inches(0.24), "%02d" % _page["n"],
    size=11, color=DIM, align=PP_ALIGN.RIGHT, spacing=1.0)

txt(s, ML, Inches(4.74), Inches(5.9), Inches(0.4),
    [[("一句话美术目标：", {"color": FN_GOLD, "bold": True}),
      ("拍一张真实的凌晨三点地铁站照片，然后只改三样东西。", {"color": ASH})]],
    size=12, spacing=1.25)
txt(s, ML, Inches(5.16), Inches(5.9), Inches(0.3), "P0 异常元素（每战场取 2–3 条，至少 1 条 P0）",
    size=11, bold=True, color=MUTED, spacing=1.0)
rows = [["站名牌", "站名逐层变成生僻字，终章层空白；拼音始终正常"],
        ["线路图", "线是闭环，没有终点站；末端站点朱红圆点正在渗色"],
        ["地面安全线", "黄线变朱红，线内侧是旧式青砖（两个年代的地面拼在一起）"],
        ["应急指示灯", "所有箭头指同一方向，且指向隧道"],
        ["乘客", "无脸；影子方向与光源不符；低头看灭屏的手机"]]
y = Inches(5.5)
for i, (a, b) in enumerate(rows):
    t = y + i * Inches(0.29)
    txt(s, ML, t, Inches(1.3), Inches(0.26), a, size=11, bold=True, color=LILAC, spacing=1.0)
    txt(s, ML + Inches(1.34), t, Inches(4.4), Inches(0.26), b, size=11, color=MUTED, spacing=1.0)

rx = Inches(6.9)
txt(s, rx, Inches(4.78), Inches(2.6), Inches(0.3), "空间类型 → 体型系统",
    size=11.5, bold=True, color=MUTED, spacing=1.0)
sp = [("站台", "开阔长条，6 格 Boss 主战场"), ("车厢内", "狭长单行道 ⭐ 巨人的舞台"),
      ("隧道 / 检修通道", "窄 + 分叉，视线阻断"), ("闸机厅", "闸机 = 可破坏掩体")]
for i, (a, b) in enumerate(sp):
    t = Inches(5.12) + i * Inches(0.42)
    txt(s, rx, t, Inches(1.6), Inches(0.26), a, size=11, bold=True,
        color=ASH if i != 1 else FN_GOLD, spacing=1.0)
    txt(s, rx, t + Inches(0.19), Inches(2.5), Inches(0.24), b, size=10, color=MUTED, spacing=1.0)

rx2 = Inches(9.75)
rect(s, rx2, Inches(4.72), Inches(2.86), Inches(1.95), fill=RGBColor(0x10, 0x14, 0x22))
rect(s, rx2, Inches(4.72), Inches(2.86), Pt(2.5), fill=LILAC)
txt(s, rx2 + Inches(0.2), Inches(4.9), Inches(2.5), Inches(0.3), "Tile 集需求",
    size=11.5, bold=True, color=LILAC, spacing=1.0)
txt(s, rx2 + Inches(0.2), Inches(5.22), Inches(2.5), Inches(1.3),
    [[("约 7 种地形 tile", {"color": ASH, "bold": True})],
     [("地砖（干/湿）· 盲道 · 站台边缘+朱红线 · 车厢地板 · 轨道 · 格栅 · 楼梯口", {})],
     [("+ 3 种危害 tile", {"color": FN_DANGER, "bold": True})],
     [("轨道 / 漏电 / 裂缝", {})]],
    size=10, color=MUTED, spacing=1.3)

# ================================================= 09 chapter 2
s = slide()
_page["n"] += 1
pic(s, "商城-百鬼夜市.png", 0, 0, W, Inches(5.1), darken=0.04, fade_bottom=0.30)
rect(s, 0, Inches(4.9), W, Inches(2.6), fill=INK)
pic(s, "商城-百鬼夜市.png", 0, 0, W, Inches(5.1), darken=0.04, fade_bottom=0.30)
rect(s, 0, 0, W, Inches(1.15), fill=INK)
pic(s, "商城-百鬼夜市.png", 0, 0, W, Inches(1.15), darken=0.60, tint=0.42, crop_top=(0.0, 0.16))
rect(s, 0, 0, Pt(4), H, fill=LANTERN)
txt(s, ML, Inches(0.34), Inches(9), Inches(0.26), "04 · 场景 · 第二章",
    size=11.5, color=LANTERN, bold=True, spacing=1.0)
txt(s, ML, Inches(0.58), Inches(9), Inches(0.45), "商场 · 百鬼市",
    size=27, bold=True, color=ASH, spacing=1.0)
txt(s, Inches(6.3), Inches(0.66), Inches(6.3), Inches(0.4),
    "「凌晨三点营业的商场，售卖着非人的商品。」", size=14.5, color=LANTERN,
    align=PP_ALIGN.RIGHT, spacing=1.0)
txt(s, Inches(12.0), Inches(6.98), Inches(0.62), Inches(0.24), "%02d" % _page["n"],
    size=11, color=DIM, align=PP_ALIGN.RIGHT, spacing=1.0)

txt(s, ML, Inches(5.02), Inches(5.9), Inches(0.44),
    [[("第一章是“空得不对”，第二章是“热闹得不对”。", {"color": ASH, "bold": True})],
     [("灯全亮、音乐在放、店都开着，但客人不是人。全游戏视觉密度最高的一章。", {"color": MUTED})]],
    size=12.5, spacing=1.3)
txt(s, ML, Inches(5.68), Inches(5.9), Inches(0.3), "P0 异常元素",
    size=11.5, bold=True, color=MUTED, spacing=1.0)
rows = [["楼层索引牌", "出现 -1F 阴货、4F、13F；楼层数不连续"],
        ["扶梯", "只上不下；或两条都在上升"],
        ["店员", "纸人：纸扎竹骨穿着正确制服，工号写在符纸上"],
        ["家电区", "纸糊的冰箱电视，带正常品牌 logo（材质错位代表作）"],
        ["价签", "单位是“年”；促销写“买一送一，寿限自理”"]]
for i, (a, b) in enumerate(rows):
    t = Inches(6.0) + i * Inches(0.28)
    txt(s, ML, t, Inches(1.3), Inches(0.26), a, size=10.5, bold=True, color=LANTERN, spacing=1.0)
    txt(s, ML + Inches(1.34), t, Inches(4.6), Inches(0.26), b, size=10.5, color=MUTED, spacing=1.0)

rx = Inches(7.1)
rect(s, rx, Inches(4.96), Inches(2.5), Inches(1.9), fill=RGBColor(0x1A, 0x0E, 0x10))
rect(s, rx, Inches(4.96), Inches(2.5), Pt(2.5), fill=FN_GOLD)
txt(s, rx + Inches(0.2), Inches(5.14), Inches(2.2), Inches(0.3), "复用红利",
    size=11.5, bold=True, color=FN_GOLD, spacing=1.0)
txt(s, rx + Inches(0.2), Inches(5.44), Inches(2.14), Inches(1.2),
    [[("货架 · 模特 · 桌椅 · 吊灯笼", {"color": ASH, "bold": True})],
     [("四类 prop 组合出全部店铺。区别靠贴图与 prop 组合，不靠新建资产。", {})]],
    size=10, color=MUTED, spacing=1.3)

rx2 = Inches(9.85)
rect(s, rx2, Inches(4.96), Inches(2.76), Inches(1.9), fill=RGBColor(0x1A, 0x0E, 0x10))
rect(s, rx2, Inches(4.96), Inches(2.76), Pt(2.5), fill=LANTERN)
txt(s, rx2 + Inches(0.2), Inches(5.14), Inches(2.4), Inches(0.3), "空间与机制",
    size=11.5, bold=True, color=LANTERN, spacing=1.0)
txt(s, rx2 + Inches(0.2), Inches(5.44), Inches(2.4), Inches(1.3),
    [[("中庭", {"color": ASH, "bold": True}), (" 6 格 Boss 主场、纵深最强", {})],
     [("扶梯 / 天井边缘", {"color": FN_GOLD, "bold": True}), (" ⭐ 推拉机制的舞台", {})],
     [("后场通道", {"color": ASH, "bold": True}), (" 冷光、无装饰 —— “后台没有伪装”", {})]],
    size=10, color=MUTED, spacing=1.28)

# ================================================= 10 chapter 3
s = slide(RGBColor(0x05, 0x0D, 0x13))
head(s, "终章 · 医院 · 阴司", kicker="04 · 场景 · 终章",
     sub="「白天是医院，夜晚是阴间的入口。」唯一一个明亮、干净、没有暗部的章节 —— 恐怖来自“过于正常”。")

grid = [("照明", "全冷白荧光，均匀铺满，几乎没有阴影"),
        ("暖色", "完全禁用。没有灯笼、没有朱红灯；唯一的红是朱批与“急”字"),
        ("材质", "全新、无垢、反光。走廊地板刚打过蜡"),
        ("异常", "更少、更精准（每战场 1–2 条），且全部是文字与数字的错误")]
for i, (a, b) in enumerate(grid):
    l = ML + (i % 2) * Inches(6.05)
    t = Inches(2.22) + (i // 2) * Inches(0.92)
    rect(s, l, t, Inches(5.65), Inches(0.76), fill=RGBColor(0x0B, 0x1E, 0x28))
    rect(s, l, t, Pt(2.5), Inches(0.76), fill=RGBColor(0x88, 0xA4, 0xB4))
    txt(s, l + Inches(0.24), t + Inches(0.1), Inches(0.9), Inches(0.3), a,
        size=13, bold=True, color=RGBColor(0x88, 0xA4, 0xB4), spacing=1.0)
    txt(s, l + Inches(1.2), t + Inches(0.09), Inches(4.28), Inches(0.62), b,
        size=10.5, color=MUTED, spacing=1.28)

txt(s, ML, Inches(4.22), Inches(11.6), Inches(0.3), "关键视觉元素：全部是文字与数字的错误",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
els = [("叫号屏", "在叫还没出生的号"), ("床位牌", "入院与“回家”是同一天"),
       ("药袋", "无药名，只有一个“忘”字"), ("病历", "朱批字迹"),
       ("太平间推床", "停在走廊中央，被单下有形状"), ("候诊椅", "坐满，全部低头，没有一个抬头"),
       ("消防疏散图", "图上的楼层平面不是这一层")]
for i, (a, b) in enumerate(els):
    l = ML + (i % 4) * Inches(2.98)
    t = Inches(4.62) + (i // 4) * Inches(0.82)
    rect(s, l, t, Inches(2.78), Inches(0.68), fill=RGBColor(0x09, 0x17, 0x1F))
    txt(s, l + Inches(0.16), t + Inches(0.08), Inches(2.5), Inches(0.26), a,
        size=11.5, bold=True, color=ASH, spacing=1.0)
    txt(s, l + Inches(0.16), t + Inches(0.35), Inches(2.5), Inches(0.26), b,
        size=10.5, color=MUTED, spacing=1.0)

rect(s, ML, Inches(6.26), CW, Inches(0.5), fill=RGBColor(0x0B, 0x1E, 0x28))
rect(s, ML, Inches(6.26), Pt(3), Inches(0.5), fill=FN_ALLY)
txt(s, ML + Inches(0.24), Inches(6.36), Inches(11.4), Inches(0.3),
    [[("资产量最省的一章", {"color": FN_ALLY, "bold": True}),
      ("（约 5 种 tile + 2 种危害）—— 干净 = 细节少。这也是把终章放在最后的另一个理由：排期最紧的时候做资产量最小的一章。", {"color": MUTED})]],
    size=12, spacing=1.0)

# ================================================= 11 four classes
s = slide()
head(s, "四职业：现代身份 × 传统力量", kicker="05 · 角色", footer=False,
     sub="每个角色都必须一眼看出“是现代人”。传统力量必须表现为“外挂”而非“自带”—— 符纸是贴上去的、法器是拿在手里的。")
pic(s, "角色.png", ML, Inches(2.08), Inches(11.89), Inches(2.62), darken=0.06,
    crop_top=(0.02, 0.62))
rect(s, ML, Inches(2.08), Inches(11.89), Inches(2.62), fill=None, line=HAIR, lw=1)

cls = [("镇邪者", "骑士 · 1 格 · 坦克", "异常事件处理员", "现代 70 / 传统 30", VERM,
        "战术背心 + 防爆盾（盾面是符板）"),
       ("方士", "法师 · 1 格 · 组合技", "民俗研究员", "现代 60 / 传统 40",
        RGBColor(0x2A, 0x6B, 0x8A), "白大褂 / 风衣 + 罗盘 · 线装书 · 符笔"),
       ("御灵者", "弓箭手 · 1 格 · 风筝", "灵异调查者", "现代 80 / 传统 20", LILAC,
        "手机 / 相机（当作法器）· 耳机 · 铃"),
       ("妖化者", "巨人 · 3 格 · 地形", "异常感染者", "现代 40 / 传统 60 · 非人 30",
        FN_DANGER, "破损的现代衣物 · 已经不是人的肢体")]
for i, (nm, proto, ident, ratio, c, gear) in enumerate(cls):
    l = ML + i * Inches(3.0)
    t = Inches(4.92)
    rect(s, l, t, Inches(2.8), Inches(1.82), fill=PANEL)
    rect(s, l, t, Inches(2.8), Pt(3), fill=c)
    txt(s, l + Inches(0.2), t + Inches(0.14), Inches(1.6), Inches(0.3), nm,
        size=16, bold=True, color=ASH, spacing=1.0)
    txt(s, l + Inches(0.2), t + Inches(0.48), Inches(2.5), Inches(0.26), proto,
        size=10, color=c if i != 1 else FN_BLOCK, bold=True, spacing=1.0)
    txt(s, l + Inches(0.2), t + Inches(0.74), Inches(2.5), Inches(0.26), ident,
        size=11, color=ASH, spacing=1.0)
    txt(s, l + Inches(0.2), t + Inches(1.0), Inches(2.5), Inches(0.26), ratio,
        size=10, color=FN_GOLD, font=MONO, spacing=1.0)
    txt(s, l + Inches(0.2), t + Inches(1.26), Inches(2.45), Inches(0.5), gear,
        size=9.5, color=MUTED, spacing=1.25)

rect(s, ML, Inches(6.86), CW, Pt(0.7), fill=HAIR)
txt(s, ML, Inches(6.94), Inches(11.0), Inches(0.3),
    [[("妖化者是唯一允许使用高饱和红的角色", {"color": FN_DANGER, "bold": True}),
      ("，因为他本身就是“危险”。这是有意设计：玩家操作他时会持续感到不安。也是唯一 3 格体型的玩家角色 —— 视觉上必须是“身体正在向外蔓延”，不是“一个很胖的人”。", {"color": MUTED})]],
    size=11, spacing=1.0)

# ================================================= 12 class sheets
s = slide()
head(s, "角色设定图：四职业已完成 L1", kicker="05 · 角色 · 现有资产",
     sub="三视图 / 武器法器 / 配饰 / 剪影方案齐备。可直接进入 sprite 拆分与骨骼动画阶段（A3：2D 手绘 sprite）。")
sheets = [("镇妖者.png", "镇邪者", "以人之力，镇压妖邪", VERM),
          ("方士.png", "方士", "以术法，改写规则", RGBColor(0x2A, 0x6B, 0x8A)),
          ("御灵者.png", "御灵者", "以怪谈，为己所用", LILAC),
          ("妖化者.png", "妖化者", "以自身，化为妖", FN_DANGER)]
for i, (f, nm, line, c) in enumerate(sheets):
    l = ML + (i % 2) * Inches(6.05)
    t = Inches(2.12) + (i // 2) * Inches(2.42)
    pic(s, f, l, t, Inches(5.65), Inches(2.18), darken=0.04)
    rect(s, l, t, Inches(5.65), Inches(2.18), fill=None, line=HAIR, lw=1)
    rect(s, l, t + Inches(1.74), Inches(2.6), Inches(0.44), fill=INK)
    rect(s, l, t + Inches(1.74), Pt(2.5), Inches(0.44), fill=c)
    txt(s, l + Inches(0.16), t + Inches(1.8), Inches(1.0), Inches(0.3), nm,
        size=13, bold=True, color=ASH, spacing=1.0)
    txt(s, l + Inches(1.1), t + Inches(1.84), Inches(1.5), Inches(0.26), line,
        size=9.5, color=MUTED, spacing=1.0)

# ================================================= 13 monsters / sizes
s = slide()
head(s, "体型系统的视觉承载：1 / 3 / 6 格", kicker="05 · 怪物",
     sub="体型系统若不可读，玩家会一直误判“我能不能走进去”。剪影语言必须随体型改变 —— 3 格不能是“放大的人形”。")

sizes = [("1 格", "直立人形，竖向剪影", "≈1.6 格", "与玩家角色同尺度，一眼是“人”", FN_ALLY, 1),
         ("3 格", "非人比例：四足 / 横向展开 / 拖行", "≈2.2 格",
          "横向剪影必须明显宽于 1 格单位", FN_GOLD, 3),
         ("6 格", "建筑 / 结构尺度，有超出屏幕感", "≈3.5 格",
          "要有明确的“结构”感（车厢、天井），不是“很大的怪”", FN_DANGER, 6)]
for i, (nm, sil, hgt, must, c, n) in enumerate(sizes):
    l = ML + i * Inches(4.02)
    rect(s, l, Inches(2.15), Inches(3.78), Inches(2.5), fill=PANEL)
    rect(s, l, Inches(2.15), Inches(3.78), Pt(3), fill=c)
    txt(s, l + Inches(0.24), Inches(2.34), Inches(1.2), Inches(0.4), nm,
        size=21, bold=True, color=c, spacing=1.0)
    txt(s, l + Inches(1.5), Inches(2.44), Inches(2.1), Inches(0.3), "高 " + hgt,
        size=11, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT, spacing=1.0)
    # footprint diagram
    hs = Inches(0.34)
    coords = {1: [(0, 0)], 3: [(0, 0), (1, 0), (0.5, 1)],
              6: [(0, 0), (1, 0), (2, 0), (0.5, 1), (1.5, 1), (1, 2)]}[n]
    rows = max(cyi for _, cyi in coords) + 1
    cx0 = l + Inches(0.46) + hs / 2
    cy0 = Inches(3.42) - (rows - 1) * hs * 0.866 / 2
    for (cxi, cyi) in coords:
        hexa(s, cx0 + cxi * hs, cy0 + cyi * hs * 0.866, hs, fill=None,
             line=c, lw=1.4)
    txt(s, l + Inches(2.0), Inches(2.9), Inches(1.62), Inches(0.72), sil,
        size=10.5, color=ASH, spacing=1.26)
    txt(s, l + Inches(0.24), Inches(4.06), Inches(3.34), Inches(0.5), must,
        size=10.5, color=MUTED, spacing=1.26)

txt(s, ML, Inches(4.9), Inches(5.6), Inches(0.3), "三条不可省略的实现要求（策划案 §13.2）",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
bullets(s, ML, Inches(5.28), Inches(5.6), [
    "footprint 完整描边 —— 所有单位按占格完整描边，不只画锚点格（Line2D 或 shader）",
    "接地阴影 —— 3/4 视角下，每个单位在每一个占用格上画一个接地阴影椭圆",
    "移动预览显示 footprint —— 拖动时实时显示落点完整占位与合法性（青蓝 / 红）",
], size=11.5, gap=0.5, spacing=1.3, marker="■", mcolor=FN_ALLY)

rx = Inches(6.9)
txt(s, rx, Inches(4.9), Inches(5.7), Inches(0.3), "两个 Boss 都是“建筑 / 载具本身”",
    size=12.5, bold=True, color=FN_DANGER, spacing=1.0)
txt(s, rx, Inches(5.28), Inches(5.7), Inches(1.4),
    [[("一章 Boss = 末班车（车本体）  ·  二章 Boss = 商场中庭（建筑本体）", {"color": ASH})],
     [("理由：① 6 格在 7×7 战场占 12%，人形放大到这个尺寸会很蠢；②“整座商场就是那个鬼”完全符合“日常物件才是异常”的调性；③ Boss 可以是战场本身的一部分，省一整套大型单位动画。", {})]],
    size=11, color=MUTED, spacing=1.34)
rect(s, rx, Inches(6.42), Inches(5.7), Inches(0.4), fill=RGBColor(0x1E, 0x0A, 0x0C))
txt(s, rx + Inches(0.2), Inches(6.5), Inches(5.4), Inches(0.3),
    [[("视觉密度与体型成反比 —— ", {"color": FN_GOLD, "bold": True}),
      ("6 格 Boss 的贴图细节要比 1 格小怪更简洁。", {"color": MUTED})]],
    size=11, spacing=1.0)

# ================================================= 14 battle readability
s = slide()
_page["n"] += 1
pic(s, "战斗场景.png", 0, 0, W, H, darken=0.42, tint=0.34)
rect(s, Inches(0.42), Inches(0.42), Inches(6.15), Inches(6.66), fill=INK)
pic(s, "战斗场景.png", Inches(6.35), Inches(0.0), Inches(6.98), H, darken=0.08,
    fade_left=0.22, focus=0.62)
txt(s, ML, Inches(0.72), Inches(5.6), Inches(0.26), "06 · 战场 · 可读性",
    size=11.5, color=FN_ALLY, bold=True, spacing=1.0)
txt(s, ML, Inches(0.96), Inches(5.6), Inches(0.5), "战场可读性硬规则",
    size=26, bold=True, color=ASH, spacing=1.0)
rect(s, ML, Inches(1.56), Inches(0.6), Pt(3.2), fill=VERM)
txt(s, ML, Inches(1.74), Inches(5.5), Inches(0.44),
    "直接服务 V3 支柱与策划案 §13.2，每条都不可协商。", size=12, color=MUTED, spacing=1.3)

rules = [("R1", "地形 tile 明度对比 ≤20%，花纹 ≤tile 面积 30%"),
         ("R2", "功能色永不出现在地形与 prop 上"),
         ("R3", "危害格必须同时用 颜色 + 图案 + 微动画 三重编码"),
         ("R4", "单位与地面的明度差 ≥40%"),
         ("R5", "每个单位在每个占用格上有接地阴影"),
         ("R6", "己方 / 敌方用轮廓光颜色区分，不靠底座圆环"),
         ("R7", "格线始终可见，最低 15% 不透明度"),
         ("R8", "所有战斗特效 ≤0.4 秒，且支持一键跳过"),
         ("R9", "特效不得遮挡目标格 —— 格心保持可见")]
for i, (code, r) in enumerate(rules):
    t = Inches(2.28) + i * Inches(0.38)
    txt(s, ML, t, Inches(0.5), Inches(0.26), code, size=11.5, bold=True,
        color=FN_ALLY, font=MONO, spacing=1.0)
    txt(s, ML + Inches(0.54), t, Inches(5.05), Inches(0.3), r, size=11,
        color=ASH, spacing=1.15)

rect(s, ML, Inches(5.86), Inches(5.5), Inches(1.02), fill=RGBColor(0x0A, 0x1F, 0x1F))
rect(s, ML, Inches(5.86), Pt(3), Inches(1.02), fill=FN_ALLY)
txt(s, ML + Inches(0.22), Inches(5.98), Inches(5.16), Inches(0.85),
    [[("R7 与氛围的冲突是最典型的取舍。", {"color": FN_ALLY, "bold": True})],
     [("“暗到看不清”恰恰是恐怖游戏想要的，但这是战棋。解法不是妥协亮度，而是把恐怖放到战斗之外：过场镜头、地图界面、事件插画可以很暗；战斗中的战场必须清晰。那 1.5 秒过场是氛围的配额，战斗回合内不再给。", {})]],
    size=10.5, color=MUTED, spacing=1.3)
txt(s, Inches(12.0), Inches(6.98), Inches(0.62), Inches(0.24), "%02d" % _page["n"],
    size=11, color=DIM, align=PP_ALIGN.RIGHT, spacing=1.0)

# ================================================= 15 tile spec
s = slide()
head(s, "六边形 Tile 技术规格与图层结构", kicker="06 · 战场 · 规格",
     sub="对齐策划案 §14.3。尖顶（pointy-top）· odd-r · y_sort_enabled = true。")

# hex diagram
cx = Inches(2.3)
cy = Inches(3.45)
hs = Inches(1.05)
for (dx, dy) in [(-0.5, -1), (0.5, -1), (-1, 0), (0, 0), (1, 0), (-0.5, 1), (0.5, 1)]:  # odd-r ring
    fill = None
    lc = STEEL
    if (dx, dy) == (0, 0):
        fill = RGBColor(0x0F, 0x2C, 0x2C)
        lc = FN_ALLY
    hexa(s, cx + dx * (hs - Inches(0.004)), cy + dy * hs * 0.862, hs,
         fill=fill, line=lc, lw=1.4)
txt(s, cx - Inches(0.52), cy - Inches(0.1), Inches(1.05), Inches(0.3), "128 px",
    size=10, color=FN_ALLY, align=PP_ALIGN.CENTER, font=MONO, spacing=1.0)
txt(s, ML, Inches(4.86), Inches(3.4), Inches(1.0),
    [[("W 128 · H 148 · 行间距 111", {"color": ASH, "bold": True, "font": MONO})],
     [("安全区 中心 96×96；单位锚点 = tile 底部中心。", {})],
     [("128×7 ≈ 896 px 战场宽，1920×1080 下留出两侧 UI 空间。", {})]],
    size=10, color=MUTED, spacing=1.3)

txt(s, Inches(4.4), Inches(2.15), Inches(3.6), Inches(0.3), "Godot 配置",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
code = ["tile_shape    = TILE_SHAPE_HEXAGON",
        "tile_layout   = TILE_LAYOUT_STACKED", "                # odd-r",
        "tile_offset_axis =", "  TILE_OFFSET_AXIS_HORIZONTAL", "                # 尖顶",
        "y_sort_enabled = true", "战场 7 × 7 = 49 格"]
rect(s, Inches(4.4), Inches(2.5), Inches(3.55), Inches(2.05), fill=RGBColor(0x08, 0x0B, 0x12),
     line=HAIR, lw=0.75)
for i, ln in enumerate(code):
    c = DIM if ln.strip().startswith("#") else (FN_ALLY if i == len(code) - 1 else MUTED)
    txt(s, Inches(4.56), Inches(2.6) + i * Inches(0.245), Inches(3.35), Inches(0.24),
        ln, size=9.5, color=c, font=MONO, spacing=1.0)

txt(s, Inches(4.4), Inches(4.72), Inches(3.6), Inches(0.3), "若采用更高分辨率",
    size=11.5, bold=True, color=FN_DANGER, spacing=1.0)
txt(s, Inches(4.4), Inches(5.04), Inches(3.55), Inches(0.5),
    "按 2 倍出图（256×296）再降采样，不要重画。", size=11, color=MUTED, spacing=1.3)

txt(s, Inches(8.35), Inches(2.15), Inches(4.2), Inches(0.3), "图层结构",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
layers = [("TerrainLayer", "静态 tile 集，每章 5–7 种", "美术", STEEL),
          ("HazardLayer", "动态 tile，需明显图案 + 微动画", "美术", FN_DANGER),
          ("HighlightLayer", "纯功能色半透明填充 + 描边，无花纹", "程序", FN_ALLY),
          ("FeatureLayer", "prop，必须与地形有明显明度差", "美术", LILAC),
          ("Node 层", "单位 + footprint 描边 + 接地阴影", "美术", ASH)]
for i, (nm, cont, who, c) in enumerate(layers):
    t = Inches(2.55) + i * Inches(0.78)
    rect(s, Inches(8.35), t, Inches(4.25), Inches(0.66), fill=PANEL)
    rect(s, Inches(8.35), t, Pt(2.5), Inches(0.66), fill=c)
    txt(s, Inches(8.5), t + Inches(0.07), Inches(2.2), Inches(0.26), nm,
        size=11, bold=True, color=c, font=MONO, spacing=1.0)
    txt(s, Inches(11.9), t + Inches(0.07), Inches(0.6), Inches(0.26), who,
        size=10, color=FN_GOLD if who == "程序" else MUTED, align=PP_ALIGN.RIGHT, spacing=1.0)
    txt(s, Inches(8.5), t + Inches(0.34), Inches(3.9), Inches(0.26), cont,
        size=10.5, color=MUTED, spacing=1.0)

txt(s, ML, Inches(5.98), Inches(7.6), Inches(0.3), "意图预告：可躲 vs 追踪（策划案 §8.7 点名，不可省略）",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
rect(s, ML, Inches(6.34), Inches(3.6), Inches(0.52), fill=PANEL)
rect(s, ML, Inches(6.34), Pt(2.5), Inches(0.52), fill=FN_DANGER)
txt(s, ML + Inches(0.18), Inches(6.42), Inches(3.3), Inches(0.44),
    [[("打击格 实线", {"color": ASH, "bold": True}), (" = 打空地 → 可躲", {"color": MUTED})]],
    size=11, spacing=1.2)
rect(s, ML + Inches(3.8), Inches(6.34), Inches(3.8), Inches(0.52), fill=PANEL)
rect(s, ML + Inches(3.8), Inches(6.34), Pt(2.5), Inches(0.52), fill=FN_DANGER)
txt(s, ML + Inches(3.98), Inches(6.42), Inches(3.5), Inches(0.44),
    [[("虚线锁定 + 连线到玩家", {"color": ASH, "bold": True}), (" = 追踪 → 躲不掉", {"color": MUTED})]],
    size=11, spacing=1.2)

# ================================================= 16 corruption
s = slide()
head(s, "腐蚀度：一张底图 + 三组贴花 + 四档后处理", kicker="07 · 腐蚀度",
     sub="本项目唯一需要美术随游戏状态动态变化的系统，也是全文性价比最高的设计 —— 玩家能“看见”自己有多贪。")

steps = [("0–2", "干净", RGBColor(0x13, 0x1A, 0x23), "基础后处理", "“这就是个地铁站”"),
         ("3–5", "渗透", RGBColor(0x2C, 0x1A, 0x24), "色相往红偏 5–8° · 雾密度 +30% · 贴花组 A（符纸、渗色、多出的影子）",
          "“刚才是不是没这些东西？”"),
         ("6–8", "侵蚀", RGBColor(0x4A, 0x18, 0x1A), "红饱和 +25% · 暗角加强 · 贴花组 B（结构扭曲、灯笼、朱红筋络） · 呼吸缩放 ±0.5%",
          "“这里已经不是地铁了”"),
         ("9+", "沦陷", RGBColor(0x7A, 0x10, 0x14), "全局偏红 · 灯光闪烁 · 贴花组 C · UI 边框渗出朱红",
          "“我该跑了”")]
for i, (rng, nm, c, how, feel) in enumerate(steps):
    l = ML + i * Inches(3.0)
    rect(s, l, Inches(2.15), Inches(2.8), Inches(2.75), fill=PANEL)
    rect(s, l, Inches(2.15), Inches(2.8), Inches(0.62), fill=c)
    txt(s, l + Inches(0.2), Inches(2.28), Inches(1.0), Inches(0.36), rng,
        size=17, bold=True, color=ASH, font=MONO, spacing=1.0)
    txt(s, l + Inches(1.5), Inches(2.34), Inches(1.1), Inches(0.3), nm,
        size=14, bold=True, color=ASH, align=PP_ALIGN.RIGHT, spacing=1.0)
    txt(s, l + Inches(0.2), Inches(2.92), Inches(2.45), Inches(1.3), how,
        size=10.5, color=MUTED, spacing=1.32)
    txt(s, l + Inches(0.2), Inches(4.42), Inches(2.45), Inches(0.4), feel,
        size=11.5, color=FN_GOLD, spacing=1.2)
    if i < 3:
        txt(s, l + Inches(2.82), Inches(3.28), Inches(0.2), Inches(0.3), "›",
            size=17, color=DIM, spacing=1.0)

txt(s, ML, Inches(5.16), Inches(7.5), Inches(0.3), "实现纪律（省钱的关键）",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
bullets(s, ML, Inches(5.52), Inches(7.4), [
    "档位切换靠后处理 + 贴花组，绝不靠重绘场景。底图一张，贴花三组，参数四档",
    "后处理集中在一个 WorldEnvironment / 全屏 shader，由 RunManager 的腐蚀度值驱动。禁止逐资产做腐蚀版本",
    "切换必须是渐变（≥1.5s），不是跳变。玩家应该“没注意到什么时候变的”",
], size=11, gap=0.52, spacing=1.26, marker="■", mcolor=FN_ALLY)

rect(s, Inches(8.4), Inches(5.16), Inches(4.2), Inches(1.6), fill=RGBColor(0x1E, 0x0A, 0x0C))
rect(s, Inches(8.4), Inches(5.16), Inches(4.2), Pt(3), fill=FN_DANGER)
txt(s, Inches(8.6), Inches(5.36), Inches(3.85), Inches(1.3),
    [[("⚠  后处理不得影响 UI 层与 HighlightLayer", {"color": FN_DANGER, "bold": True})],
     [("腐蚀度 9 时全局偏红，如果 UI 也偏红，FN_DANGER 就读不出来了 —— 这会直接摧毁红色分离。渲染顺序上 UI 必须在后处理之上。", {})]],
    size=10.5, color=MUTED, spacing=1.32)

# ================================================= 17 card / rune / gear
s = slide()
head(s, "卡牌 · 符文 · 装备：一眼分辨", kicker="07 · UI",
     sub="三者在策划案里职能不同，视觉上必须一眼分辨。符文用六边形是有意的 —— 它暗示“符文改的是这个世界的底层规则”。")

kinds = [("卡牌", "竖长方形", "符纸 / 宣纸", "按流派分色", "“我这回合要做的事”", FN_GOLD, "card"),
         ("符文", "六边形", "石 / 玉 / 骨的刻印", "紫 FN_RUNE", "“规则被改写了”", FN_RUNE, "hex"),
         ("装备", "圆角方形 / 异形", "金属 / 皮革 / 现代材质", "金 FN_GOLD + 材质本色", "“我身上带的东西”", FN_BLOCK, "gear")]
for i, (nm, shape, mat, col, one, c, kind) in enumerate(kinds):
    l = ML + i * Inches(2.72)
    rect(s, l, Inches(2.15), Inches(2.5), Inches(3.15), fill=PANEL)
    rect(s, l, Inches(2.15), Inches(2.5), Pt(3), fill=c)
    txt(s, l + Inches(0.2), Inches(2.32), Inches(1.4), Inches(0.3), nm,
        size=15, bold=True, color=ASH, spacing=1.0)
    cxx = l + Inches(1.25)
    if kind == "card":
        rect(s, cxx - Inches(0.42), Inches(2.78), Inches(0.84), Inches(1.16),
             fill=None, line=c, lw=1.5)
    elif kind == "hex":
        hexa(s, cxx, Inches(3.36), Inches(1.0), fill=None, line=c, lw=1.5)
    else:
        rect(s, cxx - Inches(0.55), Inches(2.95), Inches(1.1), Inches(0.82),
             fill=None, line=c, lw=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    for j, (k, v) in enumerate([("形状", shape), ("材质", mat), ("主色", col)]):
        t = Inches(4.06) + j * Inches(0.3)
        txt(s, l + Inches(0.2), t, Inches(0.5), Inches(0.26), k, size=10,
            color=DIM, spacing=1.0)
        txt(s, l + Inches(0.72), t, Inches(1.7), Inches(0.26), v, size=10,
            color=ASH, spacing=1.0)
    txt(s, l + Inches(0.2), Inches(5.02), Inches(2.15), Inches(0.3), one,
        size=11, color=c, spacing=1.15)

rx = Inches(8.9)
txt(s, rx, Inches(2.15), Inches(3.7), Inches(0.3), "卡牌内部信息层级",
    size=12.5, bold=True, color=FN_GOLD, spacing=1.0)
info = [("①", "体力消耗（左上角，最大的数字）", "最高频读取", FN_GOLD),
        ("②", "卡名（宋体，中）", "", ASH),
        ("③", "效果数值（等宽粗体，高亮）", "第二高频", FN_GOLD),
        ("④", "效果描述（黑体，小）", "", ASH),
        ("⑤", "插画（背景，可被文字压住）", "最低优先级", FN_DANGER),
        ("⑥", "标签（底部小字，仅检索用）", "", ASH)]
for i, (n, d, tagt, c) in enumerate(info):
    t = Inches(2.55) + i * Inches(0.4)
    txt(s, rx, t, Inches(0.3), Inches(0.26), n, size=12, color=c, bold=True, spacing=1.0)
    txt(s, rx + Inches(0.34), t, Inches(2.5), Inches(0.26), d, size=10.5,
        color=ASH, spacing=1.0)
    if tagt:
        txt(s, rx + Inches(2.9), t, Inches(1.4), Inches(0.26), tagt, size=9.5,
            color=c, align=PP_ALIGN.RIGHT, spacing=1.0)

rect(s, rx, Inches(5.1), Inches(3.7), Inches(1.32), fill=RGBColor(0x1E, 0x0A, 0x0C))
rect(s, rx, Inches(5.1), Pt(3), Inches(1.32), fill=FN_DANGER)
txt(s, rx + Inches(0.2), Inches(5.26), Inches(3.35), Inches(1.1),
    [[("⑤ 插画优先级最低，写进外包需求。", {"color": FN_DANGER, "bold": True})],
     [("策划案 §7.5 要求卡牌数值系数化（伤害随 ATK 成长），意味着卡面数值是动态的、会变成三位数 —— 插画必须给数字留出足够的低对比区域。", {})]],
    size=10.5, color=MUTED, spacing=1.32)

txt(s, ML, Inches(5.62), Inches(7.9), Inches(0.9),
    [[("符文面板（策划案 §6.5「顺序可见」）", {"color": FN_RUNE, "bold": True, "size": 12})],
     [("6 槽有序 + 顺序箭头。符文的六边形与卡牌的长方形在缩略图下绝不会混淆 —— 这是刻意的形状分工。", {"size": 11})]],
    size=11, color=MUTED, spacing=1.34)
for i in range(6):
    hexa(s, ML + Inches(0.42) + i * Inches(0.78), Inches(6.62), Inches(0.62),
         fill=RGBColor(0x1C, 0x14, 0x28), line=FN_RUNE, lw=1.2)
    if i < 5:
        txt(s, ML + Inches(0.72) + i * Inches(0.78), Inches(6.5), Inches(0.3),
            Inches(0.26), "›", size=13, color=FN_RUNE, align=PP_ALIGN.CENTER, spacing=1.0)

# ================================================= 18 red lines
s = slide()
head(s, "美术红线：明确不做什么", kicker="08 · 红线",
     sub="红线比方向更省钱 —— 它挡掉的是“画完了才发现不能用”。所有 AI prompt 必带 no gore。")

reds = ["不做血浆、内脏、断肢、尸体特写", "不做突然出现的 jump scare",
        "不做无源的紫色氛围光", "战斗中不做镜头旋转 / 抖屏遮挡",
        "场景里不出现高明度亮红（V 上限 45%）", "不把功能色用于装饰",
        "不做古风 / 仙侠画风", "不做过量东方符号堆砌（每房间 ≤3 条）",
        "不为腐蚀度重绘场景", "不做超过 0.4s 的战斗特效"]
whys = ["恐怖来自逻辑错误。血浆会拉低成廉价恐怖，且提高分级与平台风险",
        "45 分钟一局、玩几百小时。第二次就失效，第十次变成骚扰",
        "§5.1 —— 廉价感第一来源。光源必须有现实出处",
        "六边形战场旋转后玩家失去方位感；抖屏影响精确点击",
        "会摧毁“红 = 威胁”的信息通道",
        "功能色是玩家的信息通道，不是氛围工具",
        "“东方”在本项目里是符号与母题，不是画风。底子是现代写实城市",
        "堆砌 = 主题变装饰",
        "必须走后处理 + 贴花（§11）",
        "R8 —— 肉鸽玩家会玩几百小时"]
for i in range(10):
    l = ML + (i % 2) * Inches(6.05)
    t = Inches(2.15) + (i // 2) * Inches(0.86)
    rect(s, l, t, Inches(5.65), Inches(0.7), fill=PANEL)
    txt(s, l + Inches(0.16), t + Inches(0.08), Inches(0.3), Inches(0.3), "✕",
        size=13, bold=True, color=FN_DANGER, spacing=1.0)
    txt(s, l + Inches(0.52), t + Inches(0.06), Inches(5.0), Inches(0.28), reds[i],
        size=11.5, bold=True, color=ASH, spacing=1.05)
    txt(s, l + Inches(0.52), t + Inches(0.36), Inches(5.0), Inches(0.28), whys[i],
        size=10, color=MUTED, spacing=1.05)

rect(s, ML, Inches(6.5), CW, Inches(0.44), fill=RGBColor(0x1A, 0x14, 0x08))
rect(s, ML, Inches(6.5), Pt(3), Inches(0.44), fill=FN_GOLD)
txt(s, ML + Inches(0.22), Inches(6.6), Inches(11.4), Inches(0.3),
    [[("⚠ ", {"color": FN_GOLD, "bold": True}),
      ("不使用来源不明的参考图作为最终资产。moodboard 阶段的 AI 生成图与网络参考仅供内部沟通；正式资产必须原创或有明确授权。", {"color": MUTED})]],
    size=11, spacing=1.0)

# ================================================= 19 M0
s = slide()
head(s, "M0 资产清单：目标不是好看，是验证", kicker="08 · 排期",
     sub="对齐策划案 §17 的 M0（技术验证 Vertical Slice，2–3 周）。第一个 demo 纯灰盒 —— 验证“这游戏好玩吗”，不是“好看吗”。")

rows = [["MB-03 可读性板", "1 张", "完成", "前置于以下所有 UI 工作", ""],
        ["六边形 tile", "3 种", "灰盒", "验证 §9.1 尺寸参数与 odd-r 对齐", ""],
        ["高亮层", "4 种", "正式", "验证功能色与 R2", ""],
        ["单位占位", "1格×2 3格×1 6格×1", "灰盒+正式", "验证体型可读性（D8 最大风险）", "⭐"],
        ["意图图标", "6 种 + 可躲/追踪", "正式", "§8.3，策划案 §8.7 的核心验证项", "⭐"],
        ["卡牌", "边框 1 + 基石卡 3", "半正式", "验证信息层级与动态数值排版", ""],
        ["符文槽条", "6 槽 + 顺序箭头", "正式", "验证策划案 §6.5「顺序可见」", "⭐"],
        ["腐蚀度后处理", "4 档参数", "正式", "验证不影响 UI 与高亮层", "⭐"]]
cc = {}
for ri, r in enumerate(rows):
    if r[4]:
        cc[(ri, 4)] = FN_GOLD
    cc[(ri, 2)] = FN_ALLY if "正式" in r[2] else MUTED
table(s, ML, Inches(2.15), Inches(7.55),
      ["资产", "数量", "精度", "用途", ""],
      rows,
      [Inches(1.85), Inches(1.5), Inches(1.05), Inches(2.87), Inches(0.28)],
      rowh=Inches(0.38), fsize=10, cell_colors=cc)

rx = Inches(8.55)
rect(s, rx, Inches(2.15), Inches(4.06), Inches(1.5), fill=RGBColor(0x0A, 0x1F, 0x1F))
rect(s, rx, Inches(2.15), Inches(4.06), Pt(3), fill=FN_ALLY)
txt(s, rx + Inches(0.22), Inches(2.34), Inches(3.6), Inches(0.3),
    "M0 的美术验收标准只有一条", size=12, bold=True, color=FN_ALLY, spacing=1.0)
txt(s, rx + Inches(0.22), Inches(2.68), Inches(3.62), Inches(0.9),
    [[("把战斗画面缩到 25%，能分清 己方 / 敌方 / 体型 / 危害格 / 意图。", {"color": ASH})],
     [("分不清 = 参数还要调，不是“以后加特效就好了”。", {})]],
    size=11, color=MUTED, spacing=1.32)

txt(s, rx, Inches(3.86), Inches(4.0), Inches(0.3), "M1 追加（单局闭环，4–5 周）",
    size=12, bold=True, color=FN_GOLD, spacing=1.0)
bullets(s, rx, Inches(4.2), Inches(4.0), [
    "MB-01 / MB-02 完成",
    "第一章黄泉线完整 tile 集（7 地形 + 3 危害）+ prop 12 件 + 贴花组 A/B/C",
    "怪物 1 格×4、3 格×1（正式）",
    "英雄 1 个职业完整 —— 建议镇邪者（现代占比最高、最好画）",
    "卡牌边框 4 稀有度 + 20 张卡插画",
    "地图界面（盲探迷雾）+ 腐蚀度条",
], size=10, gap=0.45, spacing=1.22, marker="·", mcolor=FN_GOLD, mw=0.22)

rect(s, ML, Inches(5.42), Inches(7.55), Inches(1.32), fill=PANEL)
rect(s, ML, Inches(5.42), Pt(3), Inches(1.32), fill=VERM)
txt(s, ML + Inches(0.22), Inches(5.58), Inches(7.1), Inches(1.1),
    [[("注意：M0 没有任何美术交付物需要“画”。", {"color": ASH, "bold": True, "size": 12})],
     [("策划案 §14.3 明确 footprint 描边用 Line2D / shader 代码绘制；§12.1 纪律 3 要求战斗可无渲染跑完。上表中的“正式”项全部是功能色参数与图标规范，不是插画 —— 它们由程序 + 规范落地，不阻塞灰盒 demo。", {"size": 10.5})]],
    size=10.5, color=MUTED, spacing=1.32)

# ================================================= 20 open questions
s = slide()
head(s, "待决问题 A1–A8：确认后升 v0.2", kicker="08 · 待决",
     sub="文档其余部分已按“本文假设”推进，不等待答复。但这 8 条会影响资产量与管线，建议开工前拍板。")

qs = [("A1", "3 章 vs 5 场景", "tile 集数量，约 40% 场景资产量",
       "3 章主线 + 2 储备皮肤。要扩章优先扩无尽楼（重复即复用）"),
      ("A2", "视角确认（俯视 3/4）", "是否需要多朝向 sprite、接地阴影方案",
       "采纳 3/4。策划案 §14.3 的 y_sort 已隐含此结论"),
      ("A3", "2D 手绘 vs 像素 vs 3D", "整条管线与外包对象",
       "采纳 2D 手绘 sprite。L0 的质感只有手绘 / 写实能承接"),
      ("A4", "职业命名冲突", "所有角色资产命名、UI 文案、图鉴",
       "对外用 L0 命名，代码 id 用机制原型（hero_exorcist ↔ 骑士）"),
      ("A5", "单位朝向数量", "动画量 ×N",
       "6 朝向用 3 张图 + 水平翻转（前 / 侧 / 背），不做 6 套"),
      ("A6", "“战场即 Boss”方案", "Boss 战的实现方式",
       "强烈建议采用。省一整套大型单位动画，更符合调性"),
      ("A7", "目标分辨率与 UI 缩放基准", "全部出图尺寸",
       "基准 1920×1080，资产按 2x 出，支持 4K 降采样"),
      ("A8", "字体授权", "宋体 / 楷体的商用授权",
       "尽早确认。免费可商用备选：思源宋体、霞鹜文楷")]
y = Inches(2.12)
for i, (code, q, block, sug) in enumerate(qs):
    t = y + i * Inches(0.6)
    if i % 2 == 0:
        rect(s, ML, t, CW, Inches(0.56), fill=PANEL)
    txt(s, ML + Inches(0.16), t + Inches(0.15), Inches(0.5), Inches(0.28), code,
        size=12.5, bold=True, color=FN_GOLD, font=MONO, spacing=1.0)
    txt(s, ML + Inches(0.78), t + Inches(0.15), Inches(2.5), Inches(0.28), q,
        size=11.5, bold=True, color=ASH, spacing=1.0)
    txt(s, ML + Inches(3.4), t + Inches(0.16), Inches(3.1), Inches(0.28), block,
        size=10.5, color=MUTED, spacing=1.0)
    txt(s, ML + Inches(6.65), t + Inches(0.16), Inches(5.1), Inches(0.28), sug,
        size=10.5, color=FN_ALLY, spacing=1.0)

txt(s, ML + Inches(3.4), Inches(1.9), Inches(3.0), Inches(0.24), "阻塞什么",
    size=10, color=DIM, bold=True, spacing=1.0)
txt(s, ML + Inches(6.65), Inches(1.9), Inches(3.0), Inches(0.24), "本文建议",
    size=10, color=DIM, bold=True, spacing=1.0)

rect(s, ML, Inches(6.98), CW, Pt(0.7), fill=HAIR)

# ================================================= 21 one-pager
s = slide()
head(s, "一页速查：可打印贴墙", kicker="附录 A",
     sub="外包沟通、AI 出图、日常评审都用这一页。争议出现时，先回到这十条。", footer=False)

lines = [("定位", "90% 真实城市 + 10% 不对劲", ASH),
         ("支柱", "V3 可读性 > V1 可信度 > V2 克制    ← 冲突时永远牺牲氛围", FN_ALLY),
         ("公式", "传统母题 × 现代载体 × 一处日常逻辑错误   （每房间 ≤3 条）", LANTERN),
         ("色板", "#0B0C17 墨夜 / #2D3949 钢灰 / #3D2B45 靛紫 / #EEE7DD 香灰 / #6D0409 朱红(≤3%)", MUTED),
         ("分章", "一章 靛蓝紫(地铁)   二章 暗红橙(商场)   终章 青灰冷绿(医院·唯一明亮)", LILAC),
         ("红色", "压暗不动的红是装饰，明亮会动的红是威胁  （场景红 V≤45%，UI 红 V=100%）", FN_DANGER),
         ("高亮", "地面高亮只许两色：青蓝=合法/己方，亮红=危险/非法。其余功能色只上 UI", FN_ALLY),
         ("体型", "1格=人形 / 3格=非人比例 / 6格=建筑尺度  +  footprint描边 + 每格接地阴影", FN_GOLD),
         ("腐蚀", "一张底图 + 三组贴花 + 四档后处理   （后处理不得影响 UI 与高亮层）", FN_RUNE),
         ("红线", "无血浆 / 无 jump scare / 无无源紫光 / 场景无高饱和红 / 无古风画风", FN_DANGER),
         ("验收", "缩到 25% 还分得清 己方·敌方·体型·危害格·意图 = 合格", FN_ALLY)]
rect(s, ML, Inches(2.05), CW, Inches(4.68), fill=RGBColor(0x08, 0x0B, 0x12), line=HAIR, lw=0.75)
for i, (k, v, c) in enumerate(lines):
    t = Inches(2.24) + i * Inches(0.41)
    txt(s, ML + Inches(0.3), t, Inches(0.85), Inches(0.28), "【" + k + "】",
        size=12, bold=True, color=FN_GOLD, spacing=1.0)
    txt(s, ML + Inches(1.3), t, Inches(10.2), Inches(0.28), v, size=11.5,
        color=c, font=MONO if k in ("色板",) else SANS, spacing=1.0)

txt(s, ML, Inches(6.86), Inches(9), Inches(0.26),
    [[("v0.1 —— 与 00_系统策划案.md v0.3 对齐。A1–A8 确认后升 v0.2。", {"color": DIM})]],
    size=10, spacing=1.0)

# ---------------------------------------------------------------- save
prs.save(OUT)
shutil.rmtree(TMP, ignore_errors=True)
print("saved:", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
